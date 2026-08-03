"""
Scans a PPTX for text hidden inside <mc:AlternateContent><mc:Choice>/<mc:Fallback> wrappers.

python-pptx's slide.shapes (and anything built on it: markitdown, geometry_scan.py's
own shape walk) only recognizes direct <p:sp>/<p:pic>/<p:graphicFrame>/<p:grpSp>/<p:cxnSp>
children of the slide's spTree. When PowerPoint (commonly PowerPoint Online / recent
co-authoring saves) wraps a shape in <mc:AlternateContent> to gate an a14-namespace
feature, that shape becomes invisible to all of the above -- its text is silently
skipped, with no error and no indication anything was missed.

Usage:
    PYTHONUTF8=1 python altcontent_scan.py <file.pptx>

Run this alongside geometry_scan.py (step 3.5 of the skill) on every QA pass and
every re-check. A slide reported here means: read this text directly (this script
prints it in full) and fold it into the per-slide review like any other visible
shape -- do not treat "not found by other tools" as "not present in the file".
"""
import sys
from pptx import Presentation
from lxml import etree

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_MC = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_T = '{%s}t' % NS_A
P_CNVPR = '{%s}cNvPr' % NS_P
MC_ALT = '{%s}AlternateContent' % NS_MC
MC_CHOICE = '{%s}Choice' % NS_MC
MC_FALLBACK = '{%s}Fallback' % NS_MC


def get_text(el):
    return ''.join(t.text for t in el.iter(A_T) if t.text)


def get_name(el):
    cNvPr = el.find('.//' + P_CNVPR)
    return cNvPr.get('name') if cNvPr is not None else '(no name)'


def main(path):
    prs = Presentation(path)
    affected = []
    for i, slide in enumerate(prs.slides, start=1):
        spTree = slide.shapes._spTree
        blocks = []
        for alt in spTree.iter(MC_ALT):
            choice = alt.find(MC_CHOICE)
            src = choice if choice is not None else alt.find(MC_FALLBACK)
            label = 'Choice' if choice is not None else 'Fallback'
            if src is None:
                continue
            for shape_el in src:
                txt = get_text(shape_el).strip()
                if txt:
                    blocks.append((get_name(shape_el), label, txt))
        if blocks:
            affected.append((i, blocks))

    print(f"# סריקת mc:AlternateContent — {len(prs.slides)} שקפים")
    print()
    print(f"שקפים עם טקסט חבוי ב-mc:AlternateContent (בלתי-נראה ל-python-pptx הרגיל, "
          f"כולל markitdown/geometry_scan): **{len(affected)}**")
    if not affected:
        print("(תקין — אין תוכן חבוי)")
        return
    print()
    print("**חובה:** קרא את הטקסט המלא שלהלן ושלב אותו בבקרה הרגילה של השקף — "
          "אל תניח שמשהו לא-קיים רק כי scan.py/geometry_scan.py/markitdown לא הראו אותו.")
    print()
    for slide_num, blocks in affected:
        print(f"## שקף {slide_num} ({len(blocks)} בלוקים חבויים)")
        for name, label, txt in blocks:
            print(f"- **{name}** (mc:{label}): {txt}")
        print()


if __name__ == '__main__':
    main(sys.argv[1])
