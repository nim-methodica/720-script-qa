# -*- coding: utf-8 -*-
"""Geometry scan — reads the PPTX ITSELF (not the extracted text) so the
visually-blind checks stop guessing:

  #14  green correct-answer markers: WHERE the 00B050 mark actually sits,
       and which choice text it is nearest to
  layout: shapes that CROSS the canvas edge (partially on, partially off —
       fully-off-canvas shapes are the parked hint/feedback convention and
       are NOT flagged)
  #25/layout: text boxes whose estimated rendered text is far taller than
       the box (suspected overflow)

Also an image extractor for eye-checks (e.g. the character-select slide for
#26ב — read the two character images with vision instead of guessing from
file names):

Usage:
  PYTHONUTF8=1 python geometry_scan.py "<file.pptx>"                 # full scan
  PYTHONUTF8=1 python geometry_scan.py "<file.pptx>" --slides 4,30  # scan subset
  PYTHONUTF8=1 python geometry_scan.py "<file.pptx>" --extract-images 4 [outdir]
"""
import io
import os
import sys

from pptx import Presentation
from pptx.util import Emu

GREEN = {"00B050", "00B04F", "00B051"}   # correct-mark green (tolerance)
CROSS_TOL_IN = 0.06                       # ignore hairline overhangs
CHARS_PER_IN = 26 / 4.56                  # wrap heuristic @20pt (same as builder)
LINE_H_IN = 0.35


def _in(v):
    return Emu(v).inches


def iter_shapes(shapes, depth=0):
    for sh in shapes:
        yield sh, depth
        if sh.shape_type == 6:  # group
            for sub in iter_shapes(sh.shapes, depth + 1):
                yield sub


def shape_text(sh):
    try:
        return sh.text_frame.text.strip() if sh.has_text_frame else ""
    except Exception:
        return ""


def fill_hex(sh):
    try:
        if sh.fill.type is not None and sh.fill.type == 1:  # solid
            return str(sh.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def line_hex(sh):
    try:
        return str(sh.line.color.rgb)
    except Exception:
        return None


def scan_slide(slide, idx, sw_in, sh_in):
    out = []
    tops = []  # top-level shapes only, for edge-crossing / nearest-text
    for sh, depth in iter_shapes(slide.shapes):
        try:
            l, t = _in(sh.left), _in(sh.top)
            w, h = _in(sh.width), _in(sh.height)
        except Exception:
            continue
        if depth == 0:
            tops.append((sh, l, t, w, h))

        # green marker → nearest visible text (choice row) — #14 evidence
        fx, lx = fill_hex(sh), line_hex(sh)
        if (fx in GREEN) or (lx in GREEN):
            txt = shape_text(sh)
            near = ""
            if not txt:
                best = None
                for osh, _d in iter_shapes(slide.shapes):
                    ot = shape_text(osh)
                    if not ot or osh is sh:
                        continue
                    try:
                        d = abs(_in(osh.top) - t) + abs(_in(osh.left) - l) * 0.3
                    except Exception:
                        continue
                    if best is None or d < best[0]:
                        best = (d, ot)
                near = best[1][:60] if best else "?"
            out.append("סמן-ירוק (#14): מיקום (%.2f,%.2f)%s → הטקסט הקרוב: \"%s\"" % (
                l, t, " מילוי" if fx in GREEN else " קו-מתאר",
                (txt or near)))

        # suspected text overflow (estimated lines vs box height)
        txt = shape_text(sh)
        if txt and w > 0.5 and h > 0.2:
            cpl = max(8, int(w * CHARS_PER_IN))
            n = sum(max(1, -(-len(ln) // cpl)) for ln in txt.split("\n"))
            need = n * LINE_H_IN
            if need > h * 1.6 and n >= 4:
                out.append("חשד-גלישה: תיבה %.1f\"×%.1f\" עם ~%d שורות (\"%s…\")" % (
                    w, h, n, txt[:40]))

    # canvas-edge crossing (top-level only; fully-outside = parked, OK)
    for sh, l, t, w, h in tops:
        r, b = l + w, t + h
        inside = (r > CROSS_TOL_IN and l < sw_in - CROSS_TOL_IN and
                  b > CROSS_TOL_IN and t < sh_in - CROSS_TOL_IN)
        fully_in = (l >= -CROSS_TOL_IN and t >= -CROSS_TOL_IN and
                    r <= sw_in + CROSS_TOL_IN and b <= sh_in + CROSS_TOL_IN)
        if inside and not fully_in:
            txt = shape_text(sh)[:40]
            out.append("חוצה-גבול: (%.2f,%.2f) %.1f\"×%.1f\" \"%s\" — אלמנט שחלקו על הקנבס וחלקו בחוץ" % (
                l, t, w, h, txt))
    return out


def extract_images(path, slide_no, outdir):
    """Dump slide N's images to outdir for vision reading (e.g. #26ב)."""
    prs = Presentation(path)
    slide = prs.slides[slide_no - 1]
    os.makedirs(outdir, exist_ok=True)
    saved = []
    for sh, _d in iter_shapes(slide.shapes):
        if sh.shape_type == 13:  # picture
            img = sh.image
            name = "slide%d_%s.%s" % (slide_no, sh.shape_id, img.ext)
            p = os.path.join(outdir, name)
            with io.open(p, "wb") as f:
                f.write(img.blob)
            saved.append(p)
    print("חולצו %d תמונות משקף %d:" % (len(saved), slide_no))
    for p in saved:
        print("  " + p)


def main(argv):
    path = argv[0]
    if "--extract-images" in argv:
        i = argv.index("--extract-images")
        n = int(argv[i + 1])
        outdir = argv[i + 2] if len(argv) > i + 2 else os.path.join(
            os.environ.get("TEMP", "."), "geo_images")
        return extract_images(path, n, outdir)

    only = None
    if "--slides" in argv:
        only = {int(x) for x in argv[argv.index("--slides") + 1].split(",")}
    prs = Presentation(path)
    sw, sh_ = _in(prs.slide_width), _in(prs.slide_height)
    print("# סריקה גיאומטרית — %d שקפים (קנבס %.1f\"×%.1f\")" % (len(prs.slides), sw, sh_))
    hits = 0
    for i, slide in enumerate(prs.slides, 1):
        if only and i not in only:
            continue
        found = scan_slide(slide, i, sw, sh_)
        if found:
            hits += 1
            print("\n## שקף %d" % i)
            for f in found:
                print("- " + f)
    print("\nשקפים עם ממצא גיאומטרי: %d" % hits)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(2)
    sys.exit(main(sys.argv[1:]) or 0)
